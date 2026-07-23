"""
文档解析服务
支持 PDF、Word、Excel、PPT 格式
"""

import io
from abc import ABC, abstractmethod
from typing import Dict, List, Any, Optional
import PyPDF2
import mammoth
import openpyxl
from pptx import Presentation
from dataclasses import dataclass


@dataclass
class ParsedDocument:
    """解析后的文档"""
    content: str
    pages: int
    chunks: List[Dict[str, Any]]
    metadata: Dict[str, Any]


class DocumentParser(ABC):
    """文档解析器基类"""
    
    @abstractmethod
    async def parse(self, content: bytes, filename: str) -> ParsedDocument:
        pass
    
    @abstractmethod
    def can_parse(self, filename: str, mime_type: str) -> bool:
        pass
    
    def _chunk_text(self, text: str, chunk_size: int = 500, overlap: int = 50) -> List[Dict[str, Any]]:
        """文本分块"""
        if not text:
            return []
        
        chunks = []
        text = text.strip()
        
        for i in range(0, len(text), chunk_size - overlap):
            chunk_text = text[i:i + chunk_size]
            chunks.append({
                "text": chunk_text,
                "start": i,
                "end": min(i + chunk_size, len(text)),
                "length": len(chunk_text)
            })
        
        return chunks


class PDFParser(DocumentParser):
    """PDF解析器"""
    
    SUPPORTED_TYPES = ['application/pdf']
    SUPPORTED_EXTENSIONS = ['.pdf']
    
    def can_parse(self, filename: str, mime_type: str) -> bool:
        return (mime_type in self.SUPPORTED_TYPES or 
                any(filename.lower().endswith(ext) for ext in self.SUPPORTED_EXTENSIONS))
    
    async def parse(self, content: bytes, filename: str) -> ParsedDocument:
        try:
            reader = PyPDF2.PdfReader(io.BytesIO(content))
            text = ""
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            
            # 按段落分块（更合理的分块方式）
            paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
            
            # 合并过短的段落
            merged_chunks = []
            current_chunk = ""
            
            for para in paragraphs:
                if len(current_chunk) + len(para) < 500:
                    current_chunk += para + "\n\n"
                else:
                    if current_chunk:
                        merged_chunks.append(current_chunk.strip())
                    current_chunk = para + "\n\n"
            
            if current_chunk:
                merged_chunks.append(current_chunk.strip())
            
            chunks = [{"text": chunk, "index": i} for i, chunk in enumerate(merged_chunks)]
            
            return ParsedDocument(
                content=text,
                pages=len(reader.pages),
                chunks=chunks,
                metadata={
                    "filename": filename,
                    "pages": len(reader.pages),
                    "chunks": len(chunks)
                }
            )
        except Exception as e:
            raise ValueError(f"PDF解析失败: {str(e)}")


class WordParser(DocumentParser):
    """Word文档解析器"""
    
    SUPPORTED_TYPES = [
        'application/msword',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
    ]
    SUPPORTED_EXTENSIONS = ['.doc', '.docx']
    
    def can_parse(self, filename: str, mime_type: str) -> bool:
        return (mime_type in self.SUPPORTED_TYPES or 
                any(filename.lower().endswith(ext) for ext in self.SUPPORTED_EXTENSIONS))
    
    async def parse(self, content: bytes, filename: str) -> ParsedDocument:
        try:
            # 提取纯文本
            result = mammoth.extract_raw_text(io.BytesIO(content))
            text = result.value
            
            # 提取段落用于分块
            paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
            
            # 合并过短的段落
            merged_chunks = []
            current_chunk = ""
            
            for para in paragraphs:
                if len(current_chunk) + len(para) < 500:
                    current_chunk += para + "\n\n"
                else:
                    if current_chunk:
                        merged_chunks.append(current_chunk.strip())
                    current_chunk = para + "\n\n"
            
            if current_chunk:
                merged_chunks.append(current_chunk.strip())
            
            chunks = [{"text": chunk, "index": i} for i, chunk in enumerate(merged_chunks)]
            
            return ParsedDocument(
                content=text,
                pages=1,
                chunks=chunks,
                metadata={
                    "filename": filename,
                    "chunks": len(chunks)
                }
            )
        except Exception as e:
            raise ValueError(f"Word文档解析失败: {str(e)}")


class ExcelParser(DocumentParser):
    """Excel解析器"""
    
    SUPPORTED_TYPES = [
        'application/vnd.ms-excel',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    ]
    SUPPORTED_EXTENSIONS = ['.xls', '.xlsx']
    
    def can_parse(self, filename: str, mime_type: str) -> bool:
        return (mime_type in self.SUPPORTED_TYPES or 
                any(filename.lower().endswith(ext) for ext in self.SUPPORTED_EXTENSIONS))
    
    async def parse(self, content: bytes, filename: str) -> ParsedDocument:
        try:
            workbook = openpyxl.load_workbook(io.BytesIO(content))
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"=== 工作表: {sheet_name} ===")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([
                        str(cell) if cell is not None else ""
                        for cell in row
                        if cell is not None
                    ])
                    if row_text.strip():
                        text_parts.append(row_text)
            
            text = "\n".join(text_parts)
            chunks = self._chunk_text(text)
            
            return ParsedDocument(
                content=text,
                pages=len(workbook.sheetnames),
                chunks=chunks,
                metadata={
                    "filename": filename,
                    "sheets": workbook.sheetnames,
                    "chunks": len(chunks)
                }
            )
        except Exception as e:
            raise ValueError(f"Excel文档解析失败: {str(e)}")


class PPTParser(DocumentParser):
    """PowerPoint解析器"""
    
    SUPPORTED_TYPES = [
        'application/vnd.ms-powerpoint',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    ]
    SUPPORTED_EXTENSIONS = ['.ppt', '.pptx']
    
    def can_parse(self, filename: str, mime_type: str) -> bool:
        return (mime_type in self.SUPPORTED_TYPES or 
                any(filename.lower().endswith(ext) for ext in self.SUPPORTED_EXTENSIONS))
    
    async def parse(self, content: bytes, filename: str) -> ParsedDocument:
        try:
            presentation = Presentation(io.BytesIO(content))
            text_parts = []
            
            for i, slide in enumerate(presentation.slides):
                text_parts.append(f"=== 第{i+1}页 ===")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            
            text = "\n\n".join(text_parts)
            chunks = self._chunk_text(text)
            
            return ParsedDocument(
                content=text,
                pages=len(presentation.slides),
                chunks=chunks,
                metadata={
                    "filename": filename,
                    "slides": len(presentation.slides),
                    "chunks": len(chunks)
                }
            )
        except Exception as e:
            raise ValueError(f"PPT文档解析失败: {str(e)}")


class TextParser(DocumentParser):
    """纯文本解析器"""
    
    SUPPORTED_TYPES = ['text/plain']
    SUPPORTED_EXTENSIONS = ['.txt', '.md', '.csv', '.json']
    
    def can_parse(self, filename: str, mime_type: str) -> bool:
        return (mime_type in self.SUPPORTED_TYPES or 
                any(filename.lower().endswith(ext) for ext in self.SUPPORTED_EXTENSIONS))
    
    async def parse(self, content: bytes, filename: str) -> ParsedDocument:
        try:
            text = content.decode('utf-8')
            chunks = self._chunk_text(text)
            
            return ParsedDocument(
                content=text,
                pages=1,
                chunks=chunks,
                metadata={
                    "filename": filename,
                    "chunks": len(chunks)
                }
            )
        except UnicodeDecodeError:
            text = content.decode('gbk', errors='ignore')
            chunks = self._chunk_text(text)
            
            return ParsedDocument(
                content=text,
                pages=1,
                chunks=chunks,
                metadata={
                    "filename": filename,
                    "encoding": "gbk",
                    "chunks": len(chunks)
                }
            )


class DocumentParserFactory:
    """文档解析工厂"""
    
    def __init__(self):
        self.parsers: List[DocumentParser] = [
            PDFParser(),
            WordParser(),
            ExcelParser(),
            PPTParser(),
            TextParser()
        ]
    
    def get_parser(self, filename: str, mime_type: str) -> DocumentParser:
        for parser in self.parsers:
            if parser.can_parse(filename, mime_type):
                return parser
        raise ValueError(f"不支持的文件类型: {filename}")
    
    async def parse(self, content: bytes, filename: str, mime_type: str) -> ParsedDocument:
        parser = self.get_parser(filename, mime_type)
        return await parser.parse(content, filename)
